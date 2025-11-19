import sys
import os
# Patch sys.path for flexible script testing
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

import requests
from AppOpener import close, open as appopen
from pywhatkit import search as pywhatkit_search, playonyt
from dotenv import dotenv_values
from bs4 import BeautifulSoup
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import webbrowser
import subprocess
import keyboard
import asyncio
import logging
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import pyautogui
from Backend.utils import TempDirectoryPath
import time
import platform
from typing import Callable, Dict, List, Any, Tuple, Optional
from PIL import Image, ImageGrab
import pyperclip
from datetime import datetime
import json
from Backend.RealTimeScreenShare import (
    get_screen_analyzer,
    analyze_screen,
    find_and_click_text,
    observe_screen
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler()]
)

# Configure pyautogui safety settings
pyautogui.FAILSAFE = True
pyautogui.PAUSE = 0.1

# --------------------------------------------------------------------------- #
# Screen Capture and Real-time Screen Access
# --------------------------------------------------------------------------- #
class ScreenCapture:
    """Real-time screen capture and analysis capabilities."""
    
    @staticmethod
    def capture_screen(region: Optional[Tuple[int, int, int, int]] = None) -> Optional[Image.Image]:
        """Capture screen or specific region."""
        try:
            if region:
                x, y, width, height = region
                return ImageGrab.grab(bbox=(x, y, x + width, y + height))
            return ImageGrab.grab()
        except Exception as e:
            logging.error(f"Screen capture failed: {str(e)}")
            return None
    
    @staticmethod
    def save_screenshot(filename: Optional[str] = None, region: Optional[Tuple[int, int, int, int]] = None) -> Tuple[bool, Optional[str]]:
        """Save screenshot to file."""
        try:
            screenshot = ScreenCapture.capture_screen(region)
            if not screenshot:
                return False, "Failed to capture screen"
            
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"screenshot_{timestamp}.png"
            
            os.makedirs("data", exist_ok=True)
            filepath = os.path.join("data", filename)
            screenshot.save(filepath)
            logging.info(f"Screenshot saved to {filepath}")
            return True, filepath
        except Exception as e:
            logging.error(f"Screenshot save failed: {str(e)}")
            return False, str(e)
    
    @staticmethod
    def get_screen_size() -> Tuple[int, int]:
        """Get screen dimensions."""
        try:
            return pyautogui.size()
        except Exception:
            return (1920, 1080)  # Default fallback
    
    @staticmethod
    def get_mouse_position() -> Tuple[int, int]:
        """Get current mouse position."""
        try:
            return pyautogui.position()
        except Exception:
            return (0, 0)

# --------------------------------------------------------------------------- #
# Universal UI Navigator (Enhanced)
# --------------------------------------------------------------------------- #
SYSTEM = platform.system()
IS_MAC = SYSTEM == "Darwin"
IS_WIN = SYSTEM == "Windows"

class Navigator:
    """Enhanced UI navigator with voice-friendly controls."""

    def __init__(self, delay: float = 0.12) -> None:
        self.delay = delay

    def _do(self, action: Callable[[], Any]) -> bool:
        try:
            action()
            time.sleep(self.delay)
            return True
        except Exception as exc:
            logging.error(f"[Navigator] Action failed: {exc}")
            return False

    # Scroll & Swipe
    def scroll_up(self, amount: int = 3) -> bool:
        return self._do(lambda: pyautogui.scroll(amount))

    def scroll_down(self, amount: int = 3) -> bool:
        return self._do(lambda: pyautogui.scroll(-amount))

    def swipe_left(self, pixels: int = 200) -> bool:
        return self._do(lambda: pyautogui.drag(-pixels, 0, duration=0.4))

    def swipe_right(self, pixels: int = 200) -> bool:
        return self._do(lambda: pyautogui.drag(pixels, 0, duration=0.4))

    def swipe_up(self, pixels: int = 200) -> bool:
        return self._do(lambda: pyautogui.drag(0, -pixels, duration=0.4))

    def swipe_down(self, pixels: int = 200) -> bool:
        return self._do(lambda: pyautogui.drag(0, pixels, duration=0.4))

    # Zoom
    def zoom_in(self) -> bool:
        key = "command" if IS_MAC else "ctrl"
        return self._do(lambda: pyautogui.hotkey(key, "="))

    def zoom_out(self) -> bool:
        key = "command" if IS_MAC else "ctrl"
        return self._do(lambda: pyautogui.hotkey(key, "-"))

    def zoom_reset(self) -> bool:
        key = "command" if IS_MAC else "ctrl"
        return self._do(lambda: pyautogui.hotkey(key, "0"))

    # Basic navigation keys
    def page_up(self) -> bool:      return self._do(lambda: pyautogui.press("pageup"))
    def page_down(self) -> bool:    return self._do(lambda: pyautogui.press("pagedown"))
    def home(self) -> bool:         return self._do(lambda: pyautogui.press("home"))
    def end(self) -> bool:          return self._do(lambda: pyautogui.press("end"))
    def left(self) -> bool:         return self._do(lambda: pyautogui.press("left"))
    def right(self) -> bool:        return self._do(lambda: pyautogui.press("right"))
    def up(self) -> bool:           return self._do(lambda: pyautogui.press("up"))
    def down(self) -> bool:         return self._do(lambda: pyautogui.press("down"))
    def enter(self) -> bool:        return self._do(lambda: pyautogui.press("enter"))
    def escape(self) -> bool:       return self._do(lambda: pyautogui.press("escape"))
    def tab(self) -> bool:          return self._do(lambda: pyautogui.press("tab"))
    def backspace(self) -> bool:    return self._do(lambda: pyautogui.press("backspace"))
    def delete(self) -> bool:       return self._do(lambda: pyautogui.press("delete"))

    # Edit shortcuts (cross-platform)
    def _mod(self) -> str:
        return "command" if IS_MAC else "ctrl"

    def select_all(self) -> bool:   return self._do(lambda: pyautogui.hotkey(self._mod(), "a"))
    def copy(self) -> bool:         return self._do(lambda: pyautogui.hotkey(self._mod(), "c"))
    def paste(self) -> bool:        return self._do(lambda: pyautogui.hotkey(self._mod(), "v"))
    def cut(self) -> bool:          return self._do(lambda: pyautogui.hotkey(self._mod(), "x"))
    def undo(self) -> bool:         return self._do(lambda: pyautogui.hotkey(self._mod(), "z"))
    def redo(self) -> bool:
        if IS_MAC:
            return self._do(lambda: pyautogui.hotkey("command", "shift", "z"))
        return self._do(lambda: pyautogui.hotkey(self._mod(), "y"))

    # App-level actions
    def save(self) -> bool:         return self._do(lambda: pyautogui.hotkey(self._mod(), "s"))
    def refresh(self) -> bool:
        if IS_MAC:
            return self._do(lambda: pyautogui.hotkey("command", "r"))
        return self._do(lambda: pyautogui.press("f5"))
    
    def fullscreen(self) -> bool:   return self._do(lambda: pyautogui.press("f11"))
    
    def find(self) -> bool:         return self._do(lambda: pyautogui.hotkey(self._mod(), "f"))
    
    def close_tab(self) -> bool:    return self._do(lambda: pyautogui.hotkey(self._mod(), "w"))
    
    def new_tab(self) -> bool:      return self._do(lambda: pyautogui.hotkey(self._mod(), "t"))
    
    def next_tab(self) -> bool:     return self._do(lambda: pyautogui.hotkey(self._mod(), "tab"))
    
    def previous_tab(self) -> bool: return self._do(lambda: pyautogui.hotkey(self._mod(), "shift", "tab"))

    # PDF page jump
    def go_to_page(self, page: int) -> bool:
        def _jump():
            pyautogui.hotkey(self._mod(), "g")
            time.sleep(0.35)
            pyautogui.write(str(page))
            pyautogui.press("enter")
        return self._do(_jump)

    # Natural-language command runner
    def run(self, *commands: str) -> List[bool]:
        mapping: Dict[str, Callable[[int], bool]] = {
            # scroll
            "scroll up":     self.scroll_up,
            "scroll down":   self.scroll_down,
            # swipe
            "swipe left":    self.swipe_left,
            "swipe right":   self.swipe_right,
            "swipe up":      self.swipe_up,
            "swipe down":    self.swipe_down,
            # zoom
            "zoom in":       lambda _: self.zoom_in(),
            "zoom out":      lambda _: self.zoom_out(),
            "zoom reset":    lambda _: self.zoom_reset(),
            # page navigation
            "page up":       lambda _: self.page_up(),
            "page down":     lambda _: self.page_down(),
            "home":          lambda _: self.home(),
            "end":           lambda _: self.end(),
            # arrows
            "left":          lambda _: self.left(),
            "right":         lambda _: self.right(),
            "up":            lambda _: self.up(),
            "down":          lambda _: self.down(),
            # keys
            "enter":         lambda _: self.enter(),
            "escape":        lambda _: self.escape(),
            "tab":           lambda _: self.tab(),
            "delete":        lambda _: self.delete(),
            # edit
            "select all":    lambda _: self.select_all(),
            "copy":          lambda _: self.copy(),
            "paste":         lambda _: self.paste(),
            "cut":           lambda _: self.cut(),
            "undo":          lambda _: self.undo(),
            "redo":          lambda _: self.redo(),
            # app
            "save":          lambda _: self.save(),
            "refresh":       lambda _: self.refresh(),
            "fullscreen":   lambda _: self.fullscreen(),
            "find":          lambda _: self.find(),
            "close tab":     lambda _: self.close_tab(),
            "new tab":       lambda _: self.new_tab(),
            "next tab":      lambda _: self.next_tab(),
            "previous tab":  lambda _: self.previous_tab(),
            # pdf
            "go to page":    lambda n: self.go_to_page(n),
        }

        results: List[bool] = []
        for raw in commands:
            cmd = raw.strip().lower()
            amount = 3

            parts = cmd.split()
            if parts and parts[-1].isdigit():
                amount = int(parts[-1])
                cmd = " ".join(parts[:-1])

            action = mapping.get(cmd)
            if action:
                results.append(action(amount))
            else:
                logging.warning(f"[Navigator] Unknown command: {raw!r}")
                results.append(False)

        return results

# Global convenience instance
nav = Navigator(delay=0.12)

# Load environment variables
env_vars = {}
try:
    with open(".env", "r", encoding="utf-8") as f:
        for i, line in enumerate(f, 1):
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" not in line:
                logging.warning(f"Invalid .env line {i}: {line}")
                continue
            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip()
            if not key or not value:
                logging.warning(f"Empty key or value in .env line {i}: {line}")
                continue
            env_vars[key] = value
except Exception as e:
    logging.error(f"Failed to read .env file: {str(e)}")

GROQ_API_KEY = env_vars.get("GroqAPIKey")
EMAIL_ADDRESS = env_vars.get("EmailAddress")
EMAIL_PASSWORD = env_vars.get("EmailPassword")

# Initialize Groq client
client = None
if GROQ_API_KEY:
    try:
        client = Groq(api_key=GROQ_API_KEY)
        client.models.list()
        logging.info("Successfully initialized Groq client")
    except Exception as e:
        logging.error(f"Groq initialization failed: {str(e)}")
else:
    logging.warning("GROQ_API_KEY not found or invalid in .env file. AI features disabled.")

# Configure HTTP session
try:
    session = requests.Session()
    session.headers.update({
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    })
    logging.info("HTTP session initialized successfully")
except Exception as e:
    logging.error(f"Failed to initialize HTTP session: {str(e)}")
    session = None

# AI Configuration
SYSTEM_PROMPT = {
    "role": "system",
    "content": "You are a professional writing assistant. Generate high-quality content in clear, concise English."
}
messages = [SYSTEM_PROMPT]

# --------------------------------------------------------------------------- #
# Enhanced Automation Functions
# --------------------------------------------------------------------------- #

def GoogleSearch(query: str) -> Tuple[bool, Optional[str]]:
    """Perform Google search"""
    try:
        pywhatkit_search(query)
        logging.info(f"Google search executed for: {query}")
        return True, None
    except Exception as e:
        error_msg = f"Google search failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def Content(topic: str) -> Tuple[bool, Optional[str]]:
    """Generate and save AI content using Groq"""
    if not client:
        error_msg = "Cannot generate content: Groq client not initialized"
        logging.error(error_msg)
        return False, error_msg

    def generate_content(prompt: str) -> str:
        """Generate content using Groq"""
        try:
            messages.append({"role": "user", "content": prompt})
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages[-6:],
                temperature=0.7,
                max_tokens=2000,
                top_p=1.0
            )
            content = response.choices[0].message.content.strip()
            messages.append({"role": "assistant", "content": content})
            logging.info(f"Generated content for prompt: {prompt[:50]}...")
            return content
        except Exception as e:
            logging.error(f"Content generation failed: {str(e)}")
            return f"Error: {str(e)}"

    def save_and_open(content: str, filename: str) -> Tuple[bool, Optional[str]]:
        """Save content to file and open in editor"""
        try:
            os.makedirs("data", exist_ok=True)
            filepath = os.path.join("data", f"{filename}.txt")
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            logging.info(f"Saved content to {filepath}")
            editor = "notepad.exe" if os.name == "nt" else "gedit"
            subprocess.Popen([editor, filepath])
            return True, None
        except Exception as e:
            error_msg = f"File operation failed: {str(e)}"
            logging.error(error_msg)
            return False, error_msg

    logging.info(f"Starting content generation for: {topic}")
    clean_topic = topic.strip()[:100]
    content = generate_content(clean_topic)
    
    if content.startswith("Error:"):
        logging.error(content)
        return False, content
        
    filename = clean_topic.lower().replace(" ", "_").replace(":", "").replace("/", "_")
    success, error = save_and_open(f"Topic: {clean_topic}\n\n{content}", filename)
    return success, error

def CreateGammaPresentation(topic: str) -> Tuple[bool, Optional[str]]:
    """Generate a PowerPoint presentation using Groq"""
    if not client:
        error_msg = "Cannot generate presentation: Groq client not initialized"
        logging.error(error_msg)
        return False, error_msg

    def generate_presentation_content(prompt: str) -> str:
        """Generate presentation content using Groq"""
        try:
            messages.append({"role": "user", "content": f"Create a presentation outline about {prompt} with 6 slides: Title, Purpose, Key Features, How It Works, Example Usage, and Conclusion. For each slide, provide a title and 3-5 concise bullet points. Format as plain text with slide titles prefixed by 'Slide X: ' and bullet points prefixed by '- '. Separate slides with a blank line."})
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages[-6:],
                temperature=0.7,
                max_tokens=2000,
                top_p=1.0
            )
            content = response.choices[0].message.content.strip()
            messages.append({"role": "assistant", "content": content})
            logging.info(f"Generated presentation content for: {prompt[:50]}...")
            return content
        except Exception as e:
            logging.error(f"Presentation content generation failed: {str(e)}")
            return f"Error: {str(e)}"

    def save_and_open_presentation(content: str, filename: str) -> Tuple[bool, Optional[str]]:
        """Save presentation content to PowerPoint file and open"""
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]

            slides = content.split("\n\n")
            for slide_text in slides:
                lines = slide_text.split("\n")
                if not lines or not lines[0].startswith("Slide"):
                    continue

                slide = prs.slides.add_slide(slide_layout)
                
                title_text = lines[0].replace("Slide ", "").replace(": ", "")
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(32)

                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                text_frame = content_box.text_frame
                for bullet in lines[1:]:
                    if bullet.startswith("- "):
                        p = text_frame.add_paragraph()
                        p.text = bullet[2:].strip()
                        p.level = 0
                        p.font.size = Pt(18)
                        p.alignment = PP_ALIGN.LEFT

            os.makedirs("data", exist_ok=True)
            filepath = os.path.join("data", f"{filename}.pptx")
            prs.save(filepath)
            logging.info(f"Presentation saved to {filepath}")

            if os.name == "nt":
                os.startfile(filepath)
            else:
                subprocess.Popen(["xdg-open", filepath])
            return True, None
        except Exception as e:
            error_msg = f"Presentation file operation failed: {str(e)}"
            logging.error(error_msg)
            return False, error_msg

    logging.info(f"Starting presentation generation for: {topic}")
    clean_topic = topic.strip()[:100]
    content = generate_presentation_content(clean_topic)
    
    if content.startswith("Error:"):
        logging.error(content)
        return False, content
        
    filename = clean_topic.lower().replace(" ", "_").replace(":", "").replace("/", "_")
    success, error = save_and_open_presentation(content, filename)
    return success, error

def YoutubeSearch(query: str) -> Tuple[bool, Optional[str]]:
    """Search YouTube"""
    try:
        url = f"https://www.youtube.com/results?search_query={query.replace(' ', '+')}"
        webbrowser.open(url)
        logging.info(f"YouTube search executed for: {query}")
        return True, None
    except Exception as e:
        error_msg = f"YouTube search failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def PlayYoutube(query: str) -> Tuple[bool, Optional[str]]:
    """Play YouTube video"""
    try:
        playonyt(query)
        logging.info(f"YouTube playback started for: {query}")
        return True, None
    except Exception as e:
        error_msg = f"YouTube playback failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def OpenApp(app_name: str) -> Tuple[bool, Optional[str]]:
    """Open application with screen-aware verification"""
    try:
        appopen(app_name, match_closest=True, throw_error=True)
        logging.info(f"Opened app: {app_name}")
        
        # Wait a moment for app to open, then verify by analyzing screen
        time.sleep(1.5)
        try:
            analysis = analyze_screen()
            # Check if app window is visible (basic check)
            logging.debug(f"Screen analysis after opening {app_name}: {analysis.get('screen_size')}")
        except Exception:
            pass  # Screen analysis is optional
        
        return True, None
    except Exception:
        if not session:
            error_msg = f"Cannot open website for {app_name}: HTTP session not initialized"
            logging.error(error_msg)
            return False, error_msg
        try:
            url = f"https://www.google.com/search?q={app_name}+official+site"
            response = session.get(url)
            soup = BeautifulSoup(response.text, "html.parser")
            link = soup.find("a", {"jsname": "UWckNb"})
            if link and (href := link.get("href")):
                webbrowser.open(href)
                logging.info(f"Opened website for: {app_name}")
                return True, None
            error_msg = f"No website found for: {app_name}"
            logging.error(error_msg)
            return False, error_msg
        except Exception as e:
            error_msg = f"App opening failed: {str(e)}"
            logging.error(error_msg)
            return False, error_msg

def CloseApp(app_name: str) -> Tuple[bool, Optional[str]]:
    """Close application or folder"""
    logging.info(f"Attempting to close: {app_name}")
    try:
        if "youtube" in app_name.lower():
            try:
                close("youtube", match_closest=True, throw_error=True)
                logging.info("YouTube app closed successfully")
                return True, None
            except Exception:
                logging.info("YouTube app not found, attempting to close browser tab/process")
                try:
                    pyautogui.hotkey("ctrl", "w")
                    logging.info("Closed active browser tab")
                except Exception as e:
                    logging.error(f"Failed to close browser tab: {str(e)}")
                subprocess.run(["taskkill", "/IM", "chrome.exe", "/F"], check=False)
                subprocess.run(["taskkill", "/IM", "firefox.exe", "/F"], check=False)
                subprocess.run(["taskkill", "/IM", "edge.exe", "/F"], check=False)
                logging.info("Terminated browser processes")
                return True, None
        elif "chrome" in app_name.lower():
            subprocess.run(["taskkill", "/IM", "chrome.exe", "/F"], check=True)
            logging.info("Chrome closed successfully")
            return True, None
        elif os.path.isdir(app_name):
            try:
                subprocess.run(["taskkill", "/IM", "explorer.exe", "/F"], check=False)
                subprocess.run(["start", "explorer.exe"], shell=True, check=False)
                logging.info(f"Closed folder window: {app_name}")
                return True, None
            except Exception as e:
                error_msg = f"Failed to close folder: {str(e)}"
                logging.error(error_msg)
                return False, error_msg
        else:
            close(app_name, match_closest=True, throw_error=True)
            logging.info(f"App {app_name} closed successfully")
            return True, None
    except Exception as e:
        error_msg = f"App closing failed for {app_name}: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def System(command: str) -> Tuple[bool, Optional[str]]:
    """System controls"""
    commands_map = {
        "mute": "volume mute",
        "unmute": "volume mute",
        "volume up": "volume up",
        "volume down": "volume down",
        "shutdown": "shutdown"
    }
    
    try:
        cmd = command.lower().strip()
        if cmd == "shutdown":
            os.system("shutdown /s /t 1" if os.name == "nt" else "shutdown -h now")
            return True, None
        if cmd in commands_map:
            keyboard.press_and_release(commands_map[cmd])
            return True, None
        error_msg = f"Unknown system command: {cmd}"
        logging.warning(error_msg)
        return False, error_msg
    except Exception as e:
        error_msg = f"System command failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def SendMail(to_address: str, subject: str, body: str) -> Tuple[bool, Optional[str]]:
    """Send email"""
    try:
        if not EMAIL_ADDRESS or not EMAIL_PASSWORD:
            error_msg = "Email credentials not configured"
            logging.error(error_msg)
            return False, error_msg
            
        msg = MIMEMultipart()
        msg["From"] = EMAIL_ADDRESS
        msg["To"] = to_address
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain"))
        
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()
            server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
            server.send_message(msg)
        logging.info(f"Email sent to: {to_address}")
        return True, None
    except Exception as e:
        error_msg = f"Email sending failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

# --------------------------------------------------------------------------- #
# New Voice-Controlled Functions
# --------------------------------------------------------------------------- #

def VoiceType(text: str) -> Tuple[bool, Optional[str]]:
    """Type text using voice command with screen-aware input field detection"""
    try:
        time.sleep(0.3)  # Brief pause before typing
        
        # Optional: Verify there's a text field or input area on screen
        try:
            analysis = analyze_screen()
            if analysis.get("has_text") or analysis.get("detected_elements"):
                logging.debug("Screen contains interactive elements, proceeding with typing")
        except Exception:
            pass  # Continue even if screen analysis fails
        
        pyautogui.write(text, interval=0.05)
        logging.info(f"Voice typed text: {text[:50]}...")
        return True, None
    except Exception as e:
        error_msg = f"Voice typing failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def TakeScreenshot(filename: Optional[str] = None) -> Tuple[bool, Optional[str]]:
    """Take screenshot via voice command"""
    try:
        success, filepath = ScreenCapture.save_screenshot(filename)
        if success:
            return True, filepath
        return False, filepath
    except Exception as e:
        error_msg = f"Screenshot failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def ReadClipboard() -> Tuple[bool, Optional[str]]:
    """Read clipboard content via voice command"""
    try:
        content = pyperclip.paste()
        logging.info("Clipboard read successfully")
        return True, content
    except Exception as e:
        error_msg = f"Clipboard read failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def WriteToClipboard(text: str) -> Tuple[bool, Optional[str]]:
    """Write to clipboard via voice command"""
    try:
        pyperclip.copy(text)
        logging.info("Clipboard written successfully")
        return True, None
    except Exception as e:
        error_msg = f"Clipboard write failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def MinimizeWindow() -> Tuple[bool, Optional[str]]:
    """Minimize current window"""
    try:
        pyautogui.hotkey("win", "down") if IS_WIN else pyautogui.hotkey("command", "m")
        logging.info("Window minimized")
        return True, None
    except Exception as e:
        error_msg = f"Window minimize failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def MaximizeWindow() -> Tuple[bool, Optional[str]]:
    """Maximize current window"""
    try:
        pyautogui.hotkey("win", "up") if IS_WIN else pyautogui.hotkey("command", "ctrl", "f")
        logging.info("Window maximized")
        return True, None
    except Exception as e:
        error_msg = f"Window maximize failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def SwitchWindow() -> Tuple[bool, Optional[str]]:
    """Switch between windows"""
    try:
        pyautogui.hotkey("alt", "tab") if IS_WIN else pyautogui.hotkey("command", "tab")
        time.sleep(0.2)
        logging.info("Window switched")
        return True, None
    except Exception as e:
        error_msg = f"Window switch failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def CreateFile(filename: str, content: str = "") -> Tuple[bool, Optional[str]]:
    """Create a file via voice command"""
    try:
        os.makedirs("data", exist_ok=True)
        filepath = os.path.join("data", filename)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(content)
        logging.info(f"File created: {filepath}")
        return True, filepath
    except Exception as e:
        error_msg = f"File creation failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def ReadFile(filename: str) -> Tuple[bool, Optional[str]]:
    """Read file content via voice command"""
    try:
        filepath = os.path.join("data", filename) if not os.path.isabs(filename) else filename
        if not os.path.exists(filepath):
            error_msg = f"File not found: {filepath}"
            logging.error(error_msg)
            return False, error_msg
        with open(filepath, "r", encoding="utf-8") as f:
            content = f.read()
        logging.info(f"File read: {filepath}")
        return True, content
    except Exception as e:
        error_msg = f"File read failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def GetScreenInfo() -> Tuple[bool, Optional[str]]:
    """Get screen information (size, mouse position)"""
    try:
        width, height = ScreenCapture.get_screen_size()
        x, y = ScreenCapture.get_mouse_position()
        info = f"Screen size: {width}x{height}, Mouse position: ({x}, {y})"
        logging.info(info)
        return True, info
    except Exception as e:
        error_msg = f"Screen info failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def AnalyzeScreen() -> Tuple[bool, Optional[str]]:
    """Analyze screen content and return information about detected elements."""
    try:
        analysis = analyze_screen()
        if not analysis or "error" in analysis:
            error_msg = analysis.get("error", "Screen analysis failed") if analysis else "Screen analysis unavailable"
            return False, error_msg
        
        info_parts = []
        info_parts.append(f"Screen size: {analysis.get('screen_size', {}).get('width')}x{analysis.get('screen_size', {}).get('height')}")
        info_parts.append(f"Mouse position: {analysis.get('mouse_position')}")
        
        if analysis.get("has_text"):
            text_preview = analysis.get("text_content", "")[:100]
            info_parts.append(f"Text detected: {text_preview}...")
        
        elements = analysis.get("detected_elements", [])
        if elements:
            buttons = [e for e in elements if e.get("type") == "button_candidate"]
            if buttons:
                info_parts.append(f"Found {len(buttons)} potential buttons")
        
        info = "; ".join(info_parts)
        logging.info(f"Screen analysis: {info}")
        return True, info
    except Exception as e:
        error_msg = f"Screen analysis failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def FindTextOnScreen(search_text: str) -> Tuple[bool, Optional[str]]:
    """Find text on screen and click on it."""
    try:
        success, message = find_and_click_text(search_text)
        if success:
            logging.info(f"Found and clicked text: {search_text}")
            return True, message
        else:
            logging.warning(f"Text not found: {search_text}")
            return False, message
    except Exception as e:
        error_msg = f"Text search failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def ObserveScreen(query: str) -> Tuple[bool, Optional[str]]:
    """Observe screen and provide analysis based on query."""
    try:
        response = observe_screen(query)
        if "error" in response:
            return False, response["error"]
        
        suggestions = response.get("suggestions", [])
        if suggestions:
            result = "; ".join(suggestions)
            logging.info(f"Screen observation: {result}")
            return True, result
        else:
            analysis = response.get("screen_analysis", {})
            if analysis.get("has_text"):
                text_preview = analysis.get("text_content", "")[:100]
                return True, f"Screen contains text: {text_preview}..."
            return True, "Screen observed but no specific findings."
    except Exception as e:
        error_msg = f"Screen observation failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

def ReadScreenText() -> Tuple[bool, Optional[str]]:
    """Read all text from screen using OCR."""
    try:
        analyzer = get_screen_analyzer()
        success, text = analyzer.get_screen_text()
        if success:
            preview = text[:500] if len(text) > 500 else text
            logging.info(f"Screen text read: {len(text)} characters")
            return True, preview
        else:
            return False, text
    except Exception as e:
        error_msg = f"Screen text reading failed: {str(e)}"
        logging.error(error_msg)
        return False, error_msg

# --------------------------------------------------------------------------- #
# Command Translation and Execution
# --------------------------------------------------------------------------- #

async def TranslateAndExecute(commands: list[str]) -> Tuple[List[bool], Optional[str]]:
    """Process commands asynchronously with improved error handling"""
    tasks = []
    errors = []
    
    for command in commands:
        cmd = command.strip().lower()
        if not cmd:
            continue
            
        try:
            if cmd.startswith("open "):
                app = cmd[5:].strip()
                tasks.append(asyncio.to_thread(OpenApp, app))
            elif cmd.startswith("close "):
                app = cmd[6:].strip()
                tasks.append(asyncio.to_thread(CloseApp, app))
            elif cmd.startswith("play "):
                query = cmd[5:].strip()
                tasks.append(asyncio.to_thread(PlayYoutube, query))
            elif cmd.startswith(("write ", "content ")):
                topic = cmd.split(" ", 1)[1].strip()
                if not client:
                    errors.append(f"Skipping content generation for '{topic}': Groq client not initialized")
                    continue
                tasks.append(asyncio.to_thread(Content, topic))
            elif cmd.startswith("create presentation "):
                topic = cmd[19:].strip()
                if not client:
                    errors.append(f"Skipping presentation generation for '{topic}': Groq client not initialized")
                    continue
                tasks.append(asyncio.to_thread(CreateGammaPresentation, topic))
            elif cmd.startswith("google search "):
                query = cmd[14:].strip()
                tasks.append(asyncio.to_thread(GoogleSearch, query))
            elif cmd.startswith("youtube search "):
                query = cmd[15:].strip()
                tasks.append(asyncio.to_thread(YoutubeSearch, query))
            elif cmd.startswith("system "):
                command = cmd[7:].strip()
                tasks.append(asyncio.to_thread(System, command))
            elif cmd.startswith("send mail "):
                parts = cmd[9:].split(" about ", 1)
                to_address = parts[0].replace("to ", "").strip()
                subject_body = parts[1].split(" with ", 1) if len(parts) > 1 else ["", ""]
                tasks.append(asyncio.to_thread(
                    SendMail,
                    to_address,
                    subject_body[0].strip(),
                    subject_body[1].strip() if len(subject_body) > 1 else ""
                ))
            elif cmd.startswith("voice type ") or cmd.startswith("type "):
                text = cmd.split(" ", 2)[2].strip() if cmd.startswith("voice type ") else cmd[5:].strip()
                tasks.append(asyncio.to_thread(VoiceType, text))
            elif cmd.startswith("screenshot") or cmd.startswith("take screenshot"):
                filename = None
                if "named" in cmd or "save as" in cmd:
                    parts = cmd.split("named" if "named" in cmd else "save as", 1)
                    if len(parts) > 1:
                        filename = parts[1].strip() + ".png"
                tasks.append(asyncio.to_thread(TakeScreenshot, filename))
            elif cmd.startswith("read clipboard"):
                tasks.append(asyncio.to_thread(ReadClipboard))
            elif cmd.startswith("copy to clipboard ") or cmd.startswith("write to clipboard "):
                text = cmd.split(" ", 3)[3].strip() if "copy to clipboard" in cmd else cmd.split(" ", 3)[3].strip()
                tasks.append(asyncio.to_thread(WriteToClipboard, text))
            elif cmd.startswith("minimize window"):
                tasks.append(asyncio.to_thread(MinimizeWindow))
            elif cmd.startswith("maximize window"):
                tasks.append(asyncio.to_thread(MaximizeWindow))
            elif cmd.startswith("switch window"):
                tasks.append(asyncio.to_thread(SwitchWindow))
            elif cmd.startswith("create file "):
                parts = cmd[12:].split(" with content ", 1)
                filename = parts[0].strip()
                content = parts[1].strip() if len(parts) > 1 else ""
                tasks.append(asyncio.to_thread(CreateFile, filename, content))
            elif cmd.startswith("read file "):
                filename = cmd[10:].strip()
                tasks.append(asyncio.to_thread(ReadFile, filename))
            elif cmd.startswith("screen info") or cmd.startswith("get screen info"):
                tasks.append(asyncio.to_thread(GetScreenInfo))
            elif cmd.startswith("analyze screen") or cmd.startswith("screen analysis"):
                tasks.append(asyncio.to_thread(AnalyzeScreen))
            elif cmd.startswith("read screen") or cmd.startswith("read screen text"):
                tasks.append(asyncio.to_thread(ReadScreenText))
            elif cmd.startswith("find text ") or cmd.startswith("click on "):
                # Extract text to find
                text_to_find = cmd.split("find text ", 1)[1] if "find text" in cmd else cmd.split("click on ", 1)[1]
                tasks.append(asyncio.to_thread(FindTextOnScreen, text_to_find))
            elif cmd.startswith("observe screen") or cmd.startswith("what's on screen"):
                # Extract query if provided
                query = cmd.split("observe screen", 1)[1].strip() if "observe screen" in cmd else cmd.split("what's on screen", 1)[1].strip() if "what's on screen" in cmd else "analyze screen"
                if not query:
                    query = "analyze screen"
                tasks.append(asyncio.to_thread(ObserveScreen, query))
            else:
                logging.warning(f"Unknown command: {cmd}")
                errors.append(f"Unknown command: {cmd}")
        except Exception as e:
            error_msg = f"Command processing error: {str(e)}"
            logging.error(error_msg)
            errors.append(error_msg)

    if not tasks:
        error_msg = "No valid commands to execute" + ("; " + "; ".join(errors) if errors else "")
        return [], error_msg

    results = await asyncio.gather(*tasks, return_exceptions=True)
    
    # Process results - extract success/failure from tuples
    processed_results = []
    for r in results:
        if isinstance(r, Exception):
            processed_results.append(False)
            errors.append(str(r))
        elif isinstance(r, tuple):
            success, error = r
            processed_results.append(success)
            if error:
                errors.append(error)
        else:
            processed_results.append(bool(r))
    
    error_msg = "; ".join(errors) if errors else None
    return processed_results, error_msg

async def Automation(commands: list[str]) -> Tuple[bool, Optional[str]]:
    """Main automation controller with improved error handling"""
    logging.info("Starting automation sequence")
    results, error = await TranslateAndExecute(commands)

    if not results:
        error_msg = "No commands executed" + (f": {error}" if error else "")
        logging.error(error_msg)
        return False, error_msg

    success_count = sum(results)
    success_rate = success_count / len(results)
    logging.info(f"Automation completed with {success_rate:.0%} success rate ({success_count}/{len(results)})")
    
    if error:
        return success_count > 0, error
    return all(results), None

if __name__ == "__main__":
    # Example usage
    commands = [
        "open notepad",
        "voice type Hello, this is a test",
        "take screenshot",
        "read clipboard"
    ]

    async def main():
        result, error = await Automation(commands)
        print(f"Result: {result}, Error: {error}")

    asyncio.run(main())
